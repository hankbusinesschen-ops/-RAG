#!/usr/bin/env python3
"""
MBB-Style Presentation Generator (McKinsey / BCG / Bain conventions)
富邦金控 MA 甄選 — 年報智慧問答系統簡報

Key MBB principles applied:
  - Action Titles: every slide title is a full sentence stating the conclusion
  - Pyramid Principle: answer first, then support
  - SCR Framework: Situation → Complication → Resolution
  - BCG Color Strategy: green highlights ONE focal point, rest is neutral gray
  - Source Citations on every data slide
  - Exhibit numbering for data-heavy slides
  - Consistent slide anatomy: section tag → action title → content → source + page
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# Style Constants — BCG Green Theme (restrained palette)
# ============================================================
COLOR_PRIMARY = RGBColor(0x2D, 0x6A, 0x4F)     # Deep green (focal accents only)
COLOR_ACCENT = RGBColor(0x40, 0x91, 0x6C)       # Medium green (secondary)
COLOR_LIGHT_GREEN = RGBColor(0xD8, 0xF3, 0xDC)  # Light green tint
COLOR_GOLD = RGBColor(0xE8, 0x8D, 0x2A)         # Gold (ONE key metric only)
COLOR_DARK = RGBColor(0x1B, 0x1B, 0x1B)         # Charcoal (body, action titles)
COLOR_GRAY = RGBColor(0x6C, 0x75, 0x7D)         # Medium gray (tags, sources)
COLOR_LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)   # Card backgrounds
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_RED_LIGHT = RGBColor(0xFC, 0xE4, 0xE4)    # Risk matrix
COLOR_YELLOW_LIGHT = RGBColor(0xFE, 0xF3, 0xCD) # Risk matrix
COLOR_GREEN_LIGHT = RGBColor(0xE8, 0xF8, 0xED)  # Risk matrix
COLOR_TABLE_HEADER = RGBColor(0x2D, 0x6A, 0x4F)
COLOR_TABLE_ALT = RGBColor(0xF0, 0xF7, 0xF4)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
TOTAL_PAGES = 12

MARGIN_LEFT = Inches(0.8)
CONTENT_TOP = Inches(1.75)
CONTENT_WIDTH = Inches(11.733)


# ============================================================
# Helper Functions
# ============================================================

def set_font(run, size=11, bold=False, color=COLOR_DARK):
    """Set font properties on a run."""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"


def add_text(slide, left, top, width, height, text="", size=11,
             bold=False, color=COLOR_DARK, align=PP_ALIGN.LEFT):
    """Add a text box with single-style text."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size, bold, color)
    return box


def add_slide_header(slide, section_tag, action_title, page_num,
                     source=None, exhibit=None):
    """MBB-standard slide header: section tag + action title + footer.

    Anatomy:
      ▌ [section_tag 11pt gray]   [exhibit label, if any]
        [action_title 20pt bold dark]
      ─────────────────────────────────
      ...content zone...
      Source: xxx                                  page / 12
    """
    # Green accent bar (left edge)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(0.1), Inches(1.55)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    # Section tag (small, gray)
    tag_text = section_tag
    if exhibit:
        tag_text = f"{exhibit}  |  {section_tag}"
    add_text(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.3),
             tag_text, size=11, color=COLOR_GRAY)

    # Action title (large, bold, dark — the CONCLUSION)
    add_text(slide, Inches(0.8), Inches(0.55), Inches(11.5), Inches(0.9),
             action_title, size=20, bold=True, color=COLOR_DARK)

    # Separator line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5),
        CONTENT_WIDTH, Inches(0.015)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_LIGHT_GREEN
    line.line.fill.background()

    # Footer — source (left)
    if source:
        add_text(slide, Inches(0.8), Inches(7.05), Inches(8), Inches(0.3),
                 f"Source: {source}", size=9, color=COLOR_GRAY)

    # Footer — page number (right)
    add_text(slide, Inches(11.0), Inches(7.05), Inches(1.5), Inches(0.3),
             f"{page_num} / {TOTAL_PAGES}", size=9, color=COLOR_GRAY,
             align=PP_ALIGN.RIGHT)


def add_card(slide, left, top, width, height, title, bullets,
             fill_color=COLOR_LIGHT_GRAY, title_color=COLOR_PRIMARY,
             highlight=None):
    """Card with title, bullets, optional green highlight bar at bottom."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(0.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.1)

    # Title
    if title:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        set_font(run, size=13, bold=True, color=title_color)
        p.space_after = Pt(6)

    # Bullets
    for bullet in bullets:
        p = tf.add_paragraph()
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = f"•  {bullet}"
        set_font(run, size=11, color=COLOR_DARK)

    # Highlight bar (business value — the "so what")
    if highlight:
        p = tf.add_paragraph()
        p.space_before = Pt(10)
        run = p.add_run()
        run.text = f"→ {highlight}"
        set_font(run, size=11, bold=True, color=COLOR_PRIMARY)

    return shape


def add_styled_table(slide, left, top, width, height, headers, rows,
                     highlight_last=True):
    """BCG-styled table: green header, alternating rows, green total row."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = tbl.table

    col_width = int(width / n_cols)
    for i in range(n_cols):
        table.columns[i].width = col_width

    # Header
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_TABLE_HEADER
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                set_font(r, size=11, bold=True, color=COLOR_WHITE)

    # Data rows
    for ri, row in enumerate(rows):
        is_last = ri == len(rows) - 1
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            if highlight_last and is_last:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_PRIMARY
            elif ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_TABLE_ALT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_WHITE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    if highlight_last and is_last:
                        set_font(r, size=11, bold=True, color=COLOR_WHITE)
                    else:
                        set_font(r, size=11, bold=False, color=COLOR_DARK)

    return tbl


# ============================================================
# Slide Builders
# ============================================================

def build_cover(prs):
    """P1 — Cover."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Top accent line
    top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(0.05)
    )
    top.fill.solid()
    top.fill.fore_color.rgb = COLOR_PRIMARY
    top.line.fill.background()

    # Bottom green bar
    btm = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.6),
        SLIDE_WIDTH, Inches(0.9)
    )
    btm.fill.solid()
    btm.fill.fore_color.rgb = COLOR_PRIMARY
    btm.line.fill.background()

    # Title
    add_text(slide, Inches(1.5), Inches(1.4), Inches(10), Inches(0.9),
             "年報智慧問答系統", size=42, bold=True,
             color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)

    # Subtitle
    add_text(slide, Inches(1.5), Inches(2.4), Inches(10), Inches(0.5),
             "Retrieval-Augmented Generation for Financial Annual Reports",
             size=16, color=COLOR_GRAY, align=PP_ALIGN.CENTER)

    # Three metric badges
    badge_y = Inches(3.4)
    badge_w = Inches(3.0)
    gap = Inches(0.6)
    sx = Inches(1.8)

    metrics = [
        ("93.3%", "Overall Accuracy (28/30)", COLOR_GOLD),
        ("100%", "Basic Questions (19/19)", COLOR_GRAY),
        ("100%", "Hallucination Guard (3/3)", COLOR_GRAY),
    ]

    for i, (val, label, color) in enumerate(metrics):
        x = sx + i * (badge_w + gap)
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, badge_y,
            badge_w, Inches(1.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_LIGHT_GRAY
        bg.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        bg.line.width = Pt(0.75)

        add_text(slide, x, badge_y + Inches(0.15), badge_w, Inches(0.8),
                 val, size=44, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, x, badge_y + Inches(0.95), badge_w, Inches(0.35),
                 label, size=11, color=COLOR_GRAY, align=PP_ALIGN.CENTER)

    # Bottom bar text
    add_text(slide, Inches(1.5), Inches(6.75), Inches(10), Inches(0.3),
             "[您的姓名]   |   富邦金控數據科學組 MA 甄選   |   2026",
             size=13, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

    # Fubon label
    add_text(slide, Inches(9.5), Inches(6.1), Inches(3), Inches(0.3),
             "Fubon Financial Holdings",
             size=9, color=COLOR_GRAY, align=PP_ALIGN.RIGHT)


def build_exec_summary(prs):
    """P2 — Executive Summary (SCR format)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide,
                     "Executive Summary",
                     "以 RAG 架構將 265 頁年報轉化為即時問答，達 93.3% 準確率與 80-90% 效率提升",
                     page_num=2,
                     source="系統評估結果（2024 富邦金控年報，30 題官方測試集）")

    # Left side: SCR narrative
    scr_x = MARGIN_LEFT
    scr_w = Inches(7.5)
    y = CONTENT_TOP + Inches(0.1)

    # Situation
    add_text(slide, scr_x, y, Inches(1.5), Inches(0.3),
             "Situation", size=12, bold=True, color=COLOR_PRIMARY)
    add_text(slide, scr_x + Inches(1.6), y, Inches(5.8), Inches(0.6),
             "富邦金控年報 265+ 頁，涵蓋 10+ 子公司（人壽、銀行、證券、產險等），合併與子公司數據交織於不同章節",
             size=12, color=COLOR_DARK)

    y += Inches(0.85)
    # Complication
    add_text(slide, scr_x, y, Inches(1.5), Inches(0.3),
             "Complication", size=12, bold=True, color=COLOR_PRIMARY)
    add_text(slide, scr_x + Inches(1.6), y, Inches(5.8), Inches(0.85),
             "人工查閱單次耗時約 30 分鐘、合併報表與子公司數據易混淆（如「稅後淨利」歸屬不清）、LLM 直接回答存在幻覺風險（捏造數字或引用錯誤公司）",
             size=12, color=COLOR_DARK)

    y += Inches(1.1)
    # Resolution
    add_text(slide, scr_x, y, Inches(1.5), Inches(0.3),
             "Resolution", size=12, bold=True, color=COLOR_PRIMARY)
    add_text(slide, scr_x + Inches(1.6), y, Inches(5.8), Inches(0.85),
             "建立 RAG 年報問答系統，以三通道混合檢索 + LLM Reranking + 雙層幻覺防護達成 93.3% 準確率，並可擴展至多金控比較分析",
             size=12, color=COLOR_DARK)

    # Right side: Key metrics
    mx = Inches(8.8)
    mw = Inches(3.8)
    my = CONTENT_TOP + Inches(0.1)

    add_text(slide, mx, my, mw, Inches(0.3),
             "Key Metrics", size=12, bold=True, color=COLOR_PRIMARY)

    metrics = [
        ("93.3%", "Overall Accuracy", COLOR_GOLD, 44),
        ("100%", "Basic (19/19)", COLOR_GRAY, 32),
        ("100%", "Hallucination (3/3)", COLOR_GRAY, 32),
    ]
    for i, (val, label, color, sz) in enumerate(metrics):
        y = my + Inches(0.5) + i * Inches(1.2)
        add_text(slide, mx, y, mw, Inches(0.65),
                 val, size=sz, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, mx, y + Inches(0.6), mw, Inches(0.3),
                 label, size=10, color=COLOR_GRAY, align=PP_ALIGN.CENTER)

    # Bottom: three pillars summary
    y_bottom = CONTENT_TOP + Inches(3.6)
    pillars = [
        ("三通道混合檢索", "FAISS + BM25 + Image\nEntity-Aware Boosting"),
        ("LLM Reranking", "Gemini Cross-Encoder\n零本地記憶體"),
        ("雙層幻覺防護", "Relevance + Grounding\n競爭對手自動攔截"),
    ]
    pw = Inches(3.7)
    pg = Inches(0.3)
    for i, (title, desc) in enumerate(pillars):
        x = MARGIN_LEFT + i * (pw + pg)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y_bottom, pw, Inches(1.3)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_LIGHT_GRAY
        box.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        box.line.width = Pt(0.5)

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.1)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        set_font(run, size=12, bold=True, color=COLOR_PRIMARY)
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = desc
        set_font(run, size=10, color=COLOR_GRAY)


def build_business_context(prs):
    """P3 — Business Context."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide,
                     "業務背景",
                     "年報查閱依賴人工翻頁，每次查詢平均耗時 30 分鐘，亟需 AI 提效",
                     page_num=3,
                     source="富邦金控 113 年度股東會年報")

    # Left: Quantified pain points
    add_card(slide, MARGIN_LEFT, CONTENT_TOP, Inches(5.5), Inches(4.8),
             "現狀痛點", [
                 "265+ 頁 — 單份年報資料量",
                 "10+ 子公司 — 數據交織在不同章節",
                 "~30 分鐘 — 單次人工查詢耗時",
                 "合併 vs 子公司 — 「稅後淨利」歸屬易混淆",
                 "跨頁表格、旋轉頁面 — 傳統 OCR 無法處理",
             ], fill_color=RGBColor(0xFD, 0xF2, 0xF2),
             title_color=RGBColor(0xC0, 0x39, 0x2B))

    # Right: AI opportunity
    add_card(slide, Inches(6.8), CONTENT_TOP, Inches(5.5), Inches(4.8),
             "AI 機會", [
                 "LLM + RAG 技術成熟，可理解自然語言問題",
                 "金融年報結構固定，適合建立專屬知識庫",
                 "一次建檔、反覆查詢，邊際成本趨近於零",
                 "從人工翻頁 → AI 秒級回答",
                 "效率提升 80-90%，準確率可量化追蹤",
             ], fill_color=COLOR_GREEN_LIGHT,
             title_color=COLOR_PRIMARY)

    # Bottom callout
    callout = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_LEFT, Inches(6.35),
        CONTENT_WIDTH, Inches(0.5)
    )
    callout.fill.solid()
    callout.fill.fore_color.rgb = COLOR_PRIMARY
    callout.line.fill.background()
    tf = callout.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "目標：建立準確、可追溯、防幻覺的年報智慧問答系統"
    set_font(run, size=13, bold=True, color=COLOR_WHITE)


def build_product_showcase(prs):
    """P4 — Product Showcase / Architecture (Exhibit 1)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide,
                     "Exhibit 1  |  產品展示",
                     "七階段智慧管線將 PDF 年報轉化為可追溯的即時答案",
                     page_num=4)

    # Pipeline — 5 consolidated steps (BCG: simpler is better)
    steps = [
        "PDF 解析\nPyMuPDF + Vision",
        "TOC 分段\n+ 向量索引",
        "三通道\n混合檢索",
        "LLM\n重排序",
        "答案生成\n+ 幻覺防護",
    ]

    step_w = Inches(2.0)
    arrow_w = Inches(0.5)
    sx = Inches(0.8)
    sy = CONTENT_TOP + Inches(0.15)

    for i, label in enumerate(steps):
        x = sx + i * (step_w + arrow_w)
        # Only first and last step are dark green (BCG focal principle)
        is_focal = i == 0 or i == len(steps) - 1
        fill = COLOR_PRIMARY if is_focal else COLOR_LIGHT_GRAY
        text_c = COLOR_WHITE if is_focal else COLOR_DARK
        border = COLOR_PRIMARY if is_focal else RGBColor(0xDE, 0xDE, 0xDE)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, sy, step_w, Inches(1.1)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.color.rgb = border
        box.line.width = Pt(0.75)

        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        set_font(run, size=11, bold=True, color=text_c)

        # Gray arrow (except last)
        if i < len(steps) - 1:
            arr = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x + step_w, sy + Inches(0.3),
                arrow_w, Inches(0.5)
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = COLOR_LIGHT_GREEN
            arr.line.fill.background()

    # Three feature cards
    cy = CONTENT_TOP + Inches(1.9)
    cw = Inches(3.7)
    cg = Inches(0.3)

    features = [
        ("即時問答", [
            "自然語言提問，秒級回覆",
            "答案附頁碼來源，可追溯驗證",
            "支援查找、歸納、計算等題型",
        ]),
        ("幻覺防護", [
            "雙層檢測：相關度閾值 + 答案錨定",
            "競爭對手自動攔截",
            "風險等級指示（Low / Medium / High）",
        ]),
        ("多文件支援", [
            "同時上傳多份年報 PDF",
            "自動合併索引、標記文件來源",
            "跨年度 / 跨公司比較分析",
        ]),
    ]

    for i, (title, bullets) in enumerate(features):
        x = MARGIN_LEFT + i * (cw + cg)
        add_card(slide, x, cy, cw, Inches(2.4), title, bullets)

    # Tech stack
    add_text(slide, MARGIN_LEFT, Inches(6.35), CONTENT_WIDTH, Inches(0.3),
             "Tech:  Python  |  LangChain  |  FAISS  |  Gemini 3 Flash  |  gemini-embedding-2-preview  |  Streamlit  |  jieba",
             size=9, color=COLOR_GRAY, align=PP_ALIGN.CENTER)


def build_three_mechanisms(prs):
    """P5 — Three Key Mechanisms (Exhibit 2)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide,
                     "Exhibit 2  |  延伸設計",
                     "三項延伸機制解決金控年報的獨特挑戰，推動準確率達 93.3%",
                     page_num=5)

    cw = Inches(3.7)
    cg = Inches(0.3)

    mechanisms = [
        (
            "三通道混合檢索\n+ Entity-Aware Boosting",
            [
                "FAISS 語意 + BM25 關鍵字 + 圖片語意",
                "RRF 融合排序（K=60）",
                "實體感知：合併 ×1.5 / 非相關 ×0.7",
                "金融 jieba 詞庫（8 類 120+ 術語）",
            ],
            "正確區分「金控合併淨利」vs「子公司淨利」",
        ),
        (
            "LLM Cross-Encoder\nReranking",
            [
                "20 候選 → Gemini 排序 → Top 5",
                "零本地記憶體，純雲端 API",
                "失敗自動 fallback 原始排序",
                "比傳統 Cross-Encoder 更靈活",
            ],
            "無需 GPU，可直接部署於現有雲端",
        ),
        (
            "Dual Hallucination\nGuard",
            [
                "L1: FAISS L2 > 1.2 → 自動拒答",
                "L2: LLM 逐句驗證文件依據",
                "Scope: 國泰/中信/玉山自動攔截",
                "風險等級：Low / Medium / High",
            ],
            "100% 幻覺攔截，防止錯誤數據外流",
        ),
    ]

    for i, (title, bullets, hl) in enumerate(mechanisms):
        x = MARGIN_LEFT + i * (cw + cg)
        add_card(slide, x, CONTENT_TOP, cw, Inches(4.9),
                 title, bullets,
                 fill_color=COLOR_LIGHT_GRAY,
                 highlight=hl)


def build_results(prs):
    """P6 — Results Verification (Exhibit 3)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide,
                     "Exhibit 3  |  成果驗證",
                     "30 題官方測試達 93.3%：基本題全對、幻覺檢測零漏網",
                     page_num=6,
                     source="系統評估結果（30 題官方測試集，LLM Judge 自動評分）")

    # Left: Key metric — ONLY 93.3% is highlighted (BCG focal principle)
    add_text(slide, MARGIN_LEFT, CONTENT_TOP, Inches(3), Inches(0.8),
             "93.3%", size=52, bold=True, color=COLOR_GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, MARGIN_LEFT, CONTENT_TOP + Inches(0.75), Inches(3), Inches(0.3),
             "Overall Accuracy (28/30)", size=11, color=COLOR_GRAY, align=PP_ALIGN.CENTER)

    add_text(slide, MARGIN_LEFT, CONTENT_TOP + Inches(1.3), Inches(3), Inches(0.5),
             "100%", size=32, bold=True, color=COLOR_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, MARGIN_LEFT, CONTENT_TOP + Inches(1.75), Inches(3), Inches(0.3),
             "Basic (19/19)", size=10, color=COLOR_GRAY, align=PP_ALIGN.CENTER)

    add_text(slide, MARGIN_LEFT, CONTENT_TOP + Inches(2.3), Inches(3), Inches(0.5),
             "100%", size=32, bold=True, color=COLOR_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, MARGIN_LEFT, CONTENT_TOP + Inches(2.75), Inches(3), Inches(0.3),
             "Hallucination (3/3)", size=10, color=COLOR_GRAY, align=PP_ALIGN.CENTER)

    # Right: Table
    headers = ["類別", "題數", "答對", "準確率", "題型"]
    rows = [
        ["基本", "19", "19", "100%", "簡易查找、複合查找"],
        ["加分", "8", "6", "75%", "前後關聯、歸納資訊"],
        ["困難", "3", "3", "100%", "幻覺檢測、計算推理"],
        ["總計", "30", "28", "93.3%", ""],
    ]
    add_styled_table(slide,
                     Inches(4.0), CONTENT_TOP + Inches(0.1),
                     Inches(8.5), Inches(2.3),
                     headers, rows, highlight_last=True)

    # Example answer
    ebox = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4.0), CONTENT_TOP + Inches(2.7),
        Inches(8.5), Inches(2.8)
    )
    ebox.fill.solid()
    ebox.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    ebox.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    ebox.line.width = Pt(0.5)

    tf = ebox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.12)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "範例 — Q12 減碳比例計算（困難題）"
    set_font(run, size=11, bold=True, color=COLOR_PRIMARY)

    lines = [
        "Q：計算 2024 年度營運排放總量相較 2023 年度的減碳比例",
        "",
        "A：2023 年：8,949 + 50,091 = 59,040 噸",
        "    2024 年：4,026 + 41,568 = 45,594 噸",
        "    減碳比例 = (59,040 − 45,594) / 59,040 = 22.77%    ✓",
        "",
        "    來源：第 108、109 頁",
    ]
    for line in lines:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        set_font(run, size=10, color=COLOR_DARK)


def build_multi_doc(prs):
    """P7 — Multi-Document Query."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide,
                     "延伸功能",
                     "多文件合併查詢讓使用者直接提問跨年度問題，取代手動翻閱比對",
                     page_num=7)

    # Three-step flow
    steps = ["上傳多份\n年報 PDF", "分別建索引\n合併向量庫", "查詢結果\n自動標記來源"]
    sw = Inches(2.8)
    aw = Inches(1.2)
    sx = Inches(1.5)
    sy = CONTENT_TOP + Inches(0.1)

    for i, label in enumerate(steps):
        x = sx + i * (sw + aw)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, sy, sw, Inches(1.1)
        )
        is_focal = i == 2  # Only last step highlighted
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_PRIMARY if is_focal else COLOR_LIGHT_GRAY
        box.line.color.rgb = COLOR_PRIMARY if is_focal else RGBColor(0xDE, 0xDE, 0xDE)
        box.line.width = Pt(0.75)

        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"Step {i+1}\n{label}"
        set_font(run, size=12, bold=True,
                 color=COLOR_WHITE if is_focal else COLOR_DARK)

        if i < 2:
            arr = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, x + sw, sy + Inches(0.3),
                aw, Inches(0.5)
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = COLOR_LIGHT_GREEN
            arr.line.fill.background()

    # Scenario cards
    cy = CONTENT_TOP + Inches(1.9)
    cw = Inches(3.7)
    cg = Inches(0.3)
    scenarios = [
        ("跨年度比較", [
            "「113 vs 112 年度稅後淨利變化？」",
            "自動從兩份年報檢索對應數據",
            "結果標記文件來源與頁碼",
        ]),
        ("跨子公司分析", [
            "「各子公司 2024 風險管理策略？」",
            "Entity-Aware 辨識子公司章節",
            "彙整多個子公司對比資訊",
        ]),
        ("趨勢追蹤", [
            "「近三年 ESG 承諾變化？」",
            "多年度資料自動比對",
            "量化指標趨勢分析",
        ]),
    ]

    for i, (title, bullets) in enumerate(scenarios):
        x = MARGIN_LEFT + i * (cw + cg)
        # First card highlighted (BCG: focal point)
        fc = COLOR_LIGHT_GREEN if i == 0 else COLOR_LIGHT_GRAY
        add_card(slide, x, cy, cw, Inches(2.5), title, bullets, fill_color=fc)

    add_text(slide, MARGIN_LEFT, Inches(6.5), CONTENT_WIDTH, Inches(0.3),
             "技術：FAISS 索引合併 + Chunk ID 自動重編 + Source Metadata 追蹤 + Parent Document 展開",
             size=9, color=COLOR_GRAY, align=PP_ALIGN.CENTER)


def build_internal_apps(prs):
    """P8 — Internal Applications."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide,
                     "對內應用",
                     "系統可融入法遵、IR、風控、策略四大場景，成為集團知識基礎設施",
                     page_num=8,
                     source="業務場景分析")

    cw = Inches(5.5)
    ch = Inches(2.1)
    gx = Inches(0.5)
    gy = Inches(0.25)

    apps = [
        ("01  法遵合規", [
            "快速查找監管揭露項目",
            "董事會多元化、資本適足率即時查詢",
            "年報變更追蹤與合規檢核",
        ]),
        ("02  投資人關係 IR", [
            "股東會前快速準備 Q&A 題庫",
            "主管即時查閱數據佐證",
            "年報重點摘要自動生成",
        ]),
        ("03  風險管理", [
            "跨子公司風險類型彙整比較",
            "合規指標追蹤（CAR、逾放比等）",
            "風險事件快速定位相關揭露",
        ]),
        ("04  策略規劃", [
            "各子公司發展策略對比",
            "KPI 達成狀況查詢",
            "同業策略比較（搭配多文件功能）",
        ]),
    ]

    positions = [
        (MARGIN_LEFT, CONTENT_TOP),
        (MARGIN_LEFT + cw + gx, CONTENT_TOP),
        (MARGIN_LEFT, CONTENT_TOP + ch + gy),
        (MARGIN_LEFT + cw + gx, CONTENT_TOP + ch + gy),
    ]

    for i, ((title, bullets), (x, y)) in enumerate(zip(apps, positions)):
        add_card(slide, x, y, cw, ch, title, bullets)

    # Integration path bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_LEFT, Inches(6.3),
        CONTENT_WIDTH, Inches(0.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "整合路徑：API 串接  →  權限控管  →  部門專屬知識庫  →  集團統一知識平台"
    set_font(run, size=12, bold=True, color=COLOR_WHITE)


def build_external_extensions(prs):
    """P9 — External Extensions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide,
                     "對外延伸",
                     "將年報知識轉化為投資人服務能力：智能客服、法說會輔助、研究員查詢",
                     page_num=9,
                     source="業務場景分析")

    cw = Inches(3.7)
    cg = Inches(0.3)
    scenarios = [
        ("智能客服", [
            "投資人對話介面查詢年報數據",
            "7×24 即時回覆，降低客服負擔",
            "答案附來源頁碼，可信度高",
            "支援中文自然語言提問",
        ]),
        ("法說會輔助", [
            "自動生成 Q&A 預測題庫",
            "根據年報重點產出建議答案",
            "快速更新簡報數據佐證",
            "主管即時查核財務數據",
        ]),
        ("研究員查詢", [
            "賣方/買方研究員快速查找數據",
            "多金控比較分析",
            "特定子公司深度資料挖掘",
            "結構化數據提取與整理",
        ]),
    ]

    for i, (title, bullets) in enumerate(scenarios):
        x = MARGIN_LEFT + i * (cw + cg)
        add_card(slide, x, CONTENT_TOP, cw, Inches(3.0), title, bullets)

    # Value proposition
    vbox = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        MARGIN_LEFT, CONTENT_TOP + Inches(3.4),
        CONTENT_WIDTH, Inches(1.2)
    )
    vbox.fill.solid()
    vbox.fill.fore_color.rgb = COLOR_LIGHT_GREEN
    vbox.line.fill.background()

    tf = vbox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.12)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "價值主張"
    set_font(run, size=13, bold=True, color=COLOR_PRIMARY)

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "從「內部查詢工具」升級為「投資人服務能力」— 提升品牌專業形象，創造差異化競爭優勢"
    set_font(run, size=12, color=COLOR_DARK)

    p = tf.add_paragraph()
    p.space_before = Pt(4)
    run = p.add_run()
    run.text = "技術支撐：多文件查詢 + 實體感知 + 來源追溯 = 可信賴的對外服務基礎"
    set_font(run, size=10, color=COLOR_GRAY)


def build_risk_assessment(prs):
    """P10 — Risk Assessment."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide,
                     "風險與落地",
                     "幻覺風險透過雙層護欄控制，定位為「分析師副駕駛」而非完全取代",
                     page_num=10,
                     source="系統風險評估框架")

    # Risk matrix (left)
    mx = MARGIN_LEFT + Inches(0.6)
    my = CONTENT_TOP + Inches(0.4)
    mw = Inches(4.5)
    mh = Inches(3.2)
    half_w = int(mw / 2) - Inches(0.05)
    half_h = int(mh / 2) - Inches(0.05)
    gap = Inches(0.08)

    quadrants = [
        ("LLM 產生看似合理\n的錯誤數字", COLOR_YELLOW_LIGHT),
        ("使用者未查核\n即信任答案", COLOR_RED_LIGHT),
        ("查詢超出\n年報範圍", COLOR_GREEN_LIGHT),
        ("答案格式\n微小差異", COLOR_LIGHT_GRAY),
    ]

    positions = [
        (mx, my),
        (mx + half_w + gap, my),
        (mx, my + half_h + gap),
        (mx + half_w + gap, my + half_h + gap),
    ]

    for i, ((label, color), (x, y)) in enumerate(zip(quadrants, positions)):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, half_w, half_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(0xDE, 0xDE, 0xDE)
        shape.line.width = Pt(0.5)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.12)
        tf.margin_top = Inches(0.08)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        set_font(run, size=10, color=COLOR_DARK)

    # Axis labels
    add_text(slide, mx, my - Inches(0.25), mw, Inches(0.2),
             "← 低機率                                        高機率 →",
             size=9, color=COLOR_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, mx - Inches(0.55), my + Inches(0.2),
             Inches(0.5), Inches(1.2),
             "高\n影\n響", size=9, color=COLOR_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, mx - Inches(0.55), my + half_h + Inches(0.3),
             Inches(0.5), Inches(1.2),
             "低\n影\n響", size=9, color=COLOR_GRAY, align=PP_ALIGN.CENTER)

    # Mitigations (right)
    rx = Inches(6.5)
    rw = Inches(6.0)
    ry = CONTENT_TOP + Inches(0.1)

    add_text(slide, rx, ry, rw, Inches(0.3),
             "緩解機制", size=14, bold=True, color=COLOR_PRIMARY)

    mitigations = [
        ("雙層幻覺防護", "Relevance Threshold + Grounding Check 雙重攔截"),
        ("來源可追溯", "每個答案附頁碼，一鍵查核原文"),
        ("風險等級指示", "LOW / MEDIUM / HIGH 視覺化標示，高風險自動拒答"),
        ("競爭對手攔截", "偵測國泰/中信/玉山相關問題，自動阻擋範圍外查詢"),
        ("部署定位", "「分析師副駕駛」— 系統草擬答案，人工審核後對外使用"),
    ]

    for i, (title, desc) in enumerate(mitigations):
        y = ry + Inches(0.5) + i * Inches(0.85)

        # Dash marker with green bar
        marker = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, rx, y + Inches(0.05),
            Inches(0.06), Inches(0.25)
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = COLOR_PRIMARY
        marker.line.fill.background()

        add_text(slide, rx + Inches(0.2), y - Inches(0.02),
                 Inches(5.5), Inches(0.25),
                 title, size=12, bold=True, color=COLOR_DARK)
        add_text(slide, rx + Inches(0.2), y + Inches(0.25),
                 Inches(5.5), Inches(0.3),
                 desc, size=10, color=COLOR_GRAY)


def build_future_roadmap(prs):
    """P11 — Future Roadmap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide,
                     "未來方向",
                     "三階段演進：API 服務化 → Agent 自動化 → 主動知識服務",
                     page_num=11)

    # Timeline bar
    bar_y = CONTENT_TOP + Inches(0.15)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.9), bar_y,
        Inches(11.4), Inches(0.06)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_LIGHT_GREEN
    bar.line.fill.background()

    phases = [
        ("Phase 1  短期", "API 服務化", [
            "REST API 串接內部工具與 BI 平台",
            "排程批次處理分析師問題集",
            "新年報上架自動重建索引",
            "權限控管與使用量追蹤",
        ], COLOR_PRIMARY),
        ("Phase 2  中期", "Agent 自動化", [
            "複雜問題自動拆解、多步推理",
            "整合計算工具（財務比率、圖表）",
            "串接外部資料（Bloomberg、CRM）",
            "Multi-turn 對話記憶",
        ], COLOR_ACCENT),
        ("Phase 3  長期", "主動知識服務", [
            "監管申報變更自動警示",
            "自動產出同業比較報告",
            "知識圖譜建構與實體關聯",
            "推播式年報重大變化通知",
        ], RGBColor(0x6C, 0xAE, 0x75)),
    ]

    pw = Inches(3.6)
    pg = Inches(0.35)
    psx = Inches(0.9)

    for i, (phase_label, phase_title, bullets, color) in enumerate(phases):
        x = psx + i * (pw + pg)

        # Timeline dot
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x + int(pw / 2) - Inches(0.15), bar_y - Inches(0.09),
            Inches(0.25), Inches(0.25)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()

        # Phase header — gray bg with green left bar (BCG section marker)
        header_y = CONTENT_TOP + Inches(0.55)
        # Left green bar on header
        left_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, header_y,
            Inches(0.08), Inches(0.65)
        )
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = color
        left_bar.line.fill.background()

        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x + Inches(0.08), header_y,
            pw - Inches(0.08), Inches(0.65)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = COLOR_LIGHT_GRAY
        header.line.fill.background()
        tf = header.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.12)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"{phase_label}  —  {phase_title}"
        set_font(run, size=12, bold=True, color=COLOR_DARK)

        # Bullet card
        add_card(slide, x, header_y + Inches(0.8),
                 pw, Inches(3.0), "", bullets)

    # Maturity arrow
    add_text(slide, Inches(0.9), Inches(6.35), Inches(11.4), Inches(0.3),
             "被動查詢  ─────────────────────────→  主動知識服務",
             size=12, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)


def build_qa(prs):
    """P12 — Q&A."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Top accent
    top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(0.05)
    )
    top.fill.solid()
    top.fill.fore_color.rgb = COLOR_PRIMARY
    top.line.fill.background()

    # Q&A
    add_text(slide, Inches(2), Inches(2.0), Inches(9), Inches(1.2),
             "Q & A", size=72, bold=True, color=COLOR_PRIMARY,
             align=PP_ALIGN.CENTER)

    add_text(slide, Inches(2), Inches(3.5), Inches(9), Inches(0.5),
             "Thank You", size=22, color=COLOR_GRAY,
             align=PP_ALIGN.CENTER)

    # Separator
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.3),
        Inches(2.3), Inches(0.02)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = COLOR_LIGHT_GREEN
    sep.line.fill.background()

    add_text(slide, Inches(2), Inches(4.7), Inches(9), Inches(0.4),
             "[您的姓名]   |   富邦金控數據科學組 MA 甄選   |   2026",
             size=13, color=COLOR_GRAY, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(2), Inches(5.5), Inches(9), Inches(0.8),
             "延伸討論：如何應用於即時法說會逐字稿？多語言年報擴展？與現有文管系統整合？",
             size=11, color=COLOR_GRAY, align=PP_ALIGN.CENTER)

    # Bottom bar
    btm = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.6),
        SLIDE_WIDTH, Inches(0.9)
    )
    btm.fill.solid()
    btm.fill.fore_color.rgb = COLOR_PRIMARY
    btm.line.fill.background()


# ============================================================
# Main
# ============================================================

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    builders = [
        build_cover,              # P1
        build_exec_summary,       # P2
        build_business_context,   # P3
        build_product_showcase,   # P4
        build_three_mechanisms,   # P5
        build_results,            # P6
        build_multi_doc,          # P7
        build_internal_apps,      # P8
        build_external_extensions,# P9
        build_risk_assessment,    # P10
        build_future_roadmap,     # P11
        build_qa,                 # P12
    ]

    for builder in builders:
        builder(prs)

    output = "presentation.pptx"
    prs.save(output)
    print(f"Saved: {output}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
